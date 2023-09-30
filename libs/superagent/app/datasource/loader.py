import json
import tempfile
from tempfile import NamedTemporaryFile
from typing import Any
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup as Soup
from langchain.docstore.document import Document
from langchain.document_loaders import (
    GitLoader,
    PyPDFLoader,
    RecursiveUrlLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
    UnstructuredWordDocumentLoader,
    WebBaseLoader,
    YoutubeLoader,
)
from langchain.document_loaders.airbyte import AirbyteStripeLoader
from llama_index import download_loader
from pyairtable import Api

from prisma.models import Datasource


class DataLoader:
    def __init__(self, datasource: Datasource):
        self.datasource = datasource

    def load(self) -> Any:
        if self.datasource.type == "TXT":
            return self.load_txt()
        elif self.datasource.type == "PDF":
            return self.load_pdf()
        elif self.datasource.type == "PPTX":
            return self.load_pptx()
        elif self.datasource.type == "DOCX":
            return self.load_docx()
        elif self.datasource.type == "GOOGLE_DOC":
            return self.load_google_doc()
        elif self.datasource.type == "Markdown":
            return self.load_markdown()
        elif self.datasource.type == "GITHUB_REPOSITORY":
            return self.load_github()
        elif self.datasource.type == "WEBPAGE":
            return self.load_webpage()
        elif self.datasource.type == "NOTION":
            return self.load_notion()
        elif self.datasource.type == "YOUTUBE":
            return self.load_youtube()
        elif self.datasource.type == "URL":
            return self.load_url()
        elif self.datasource.type == "AIRTABLE":
            return self.load_airtable()
        elif self.datasource.type == "STRIPE":
            return self.load_stripe()
        else:
            raise ValueError(f"Unsupported datasource type: {self.datasource.type}")

    def load_txt(self):
        file_response = requests.get(self.datasource.url).text
        with NamedTemporaryFile(suffix=".txt", delete=True) as temp_file:
            temp_file.write(file_response.encode())
            temp_file.flush()
            loader = TextLoader(file_path=temp_file.name)
            return loader.load_and_split()

    def load_pdf(self):
        loader = PyPDFLoader(file_path=self.datasource.url)
        return loader.load_and_split()

    def load_google_doc(self):
        pass

    def load_pptx(self):
        from pptx import Presentation

        file_response = requests.get(self.datasource.url).content
        with NamedTemporaryFile(suffix=".pptx", delete=True) as temp_file:
            temp_file.write(file_response)
            temp_file.flush()
            presentation = Presentation(temp_file.name)
            result = ""
            for i, slide in enumerate(presentation.slides):
                result += f"\n\nSlide #{i}: \n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += f"{shape.text}\n"
            return [Document(page_content=result)]

    def load_docx(self):
        file_response = requests.get(self.datasource.url).content
        with NamedTemporaryFile(suffix=".docx", delete=True) as temp_file:
            temp_file.write(file_response)
            temp_file.flush()
            loader = UnstructuredWordDocumentLoader(file_path=temp_file.name)
            return loader.load_and_split()

    def load_markdown(self):
        if file_response := requests.get(self.datasource.url).text:
            with NamedTemporaryFile(suffix=".md", delete=True) as temp_file:
                temp_file.write(file_response.encode())
                temp_file.flush()
                loader = UnstructuredMarkdownLoader(file_path=temp_file.name)
                return loader.load()

    def load_github(self):
        parsed_url = urlparse(self.datasource.url)
        path_parts = parsed_url.path.split("/")  # type: ignore
        repo_name = path_parts[2]
        metadata = json.loads(self.datasource.metadata)

        with tempfile.TemporaryDirectory() as temp_dir:
            repo_path = f"{temp_dir}/{repo_name}/"  # type: ignore
            loader = GitLoader(
                clone_url=self.datasource.url,
                repo_path=repo_path,
                branch=metadata["branch"],  # type: ignore
            )
            return loader.load_and_split()

    def load_webpage(self):
        loader = RecursiveUrlLoader(
            url=self.datasource.url,
            max_depth=2,
            extractor=lambda x: Soup(x, "html.parser").text,
        )
        chunks = loader.load_and_split()
        for chunk in chunks:
            if "language" in chunk.metadata:
                del chunk.metadata["language"]
        return chunks

    def load_notion(self):
        metadata = json.loads(self.datasource.metadata)
        NotionPageReader = download_loader("NotionPageReader")
        integration_token = metadata["integration_token"]
        page_ids = metadata["page_ids"]
        loader = NotionPageReader(integration_token=integration_token)
        return loader.load_langchain_documents(page_ids=page_ids.split(","))

    def load_youtube(self):
        video_id = self.datasource.url.split("youtube.com/watch?v=")[-1]
        loader = YoutubeLoader(video_id=video_id)
        return loader.load_and_split()

    def load_url(self):
        url_list = self.datasource.url.split(",")
        loader = WebBaseLoader(url_list)
        return loader.load_and_split()

    def load_airtable(self):
        metadata = json.loads(self.datasource.metadata)
        api_key = metadata["apiKey"]
        base_id = metadata["baseId"]
        table_id = metadata["tableId"]
        api = Api(api_key)
        table = api.table(base_id, table_id)
        return table.all()

    def load_stripe(self):
        metadata = json.loads(self.datasource.metadata)
        client_secret = metadata["clientSecret"]
        account_id = metadata["accountId"]
        start_date = metadata["startDate"]
        stream_name = metadata["streamName"]
        config = {
            "client_secret": client_secret,
            "account_id": account_id,
            "start_date": start_date,
        }

        def handle_record(record: dict, _id: str):
            return record.data

        loader = AirbyteStripeLoader(
            config=config,
            record_handler=handle_record,
            stream_name=stream_name,
        )
        data = loader.load()
        return data
